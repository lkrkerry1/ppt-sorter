"""
PPT解析器 - 统一处理PPT文件
"""

import re
import zipfile
from typing import List, Optional, Union, Any
import warnings


class PPTParser:
    """PPT解析器"""

    def __init__(self, use_fast_mode: bool = True):
        """
        初始化解析器

        Args:
            use_fast_mode: 是否使用快速模式（牺牲一些准确性）
        """
        self.use_fast_mode = use_fast_mode

    def extract_text(self, file_path: str, max_slides: int = 50) -> str:
        """
        从PPT提取文本

        Args:
            file_path: PPT文件路径
            max_slides: 最大处理幻灯片数量

        Returns:
            提取的文本
        """
        if file_path.lower().endswith(".pptx"):
            if self.use_fast_mode:
                return self._extract_from_pptx_fast(file_path, max_slides)
            else:
                return self._extract_from_pptx_full(file_path, max_slides)
        elif file_path.lower().endswith(".ppt"):
            return self._extract_from_ppt_old(file_path, max_slides)
        else:
            raise ValueError(f"不支持的文件格式: {file_path}")

    def _extract_from_pptx_fast(self, file_path: str, max_slides: int) -> str:
        """快速提取.pptx文本（直接解析XML）"""
        try:
            texts = []
            with zipfile.ZipFile(file_path, "r") as zip_ref:
                # 获取幻灯片文件
                slide_files = []
                for name in zip_ref.namelist():
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                        try:
                            # 提取幻灯片编号
                            parts = name.split("slide")
                            if len(parts) > 1:
                                slide_num_part = parts[1].split(".")[0]
                                slide_num = (
                                    int(slide_num_part)
                                    if slide_num_part.isdigit()
                                    else 0
                                )
                                slide_files.append((slide_num, name))
                        except:
                            continue

                # 排序并限制数量
                slide_files.sort(key=lambda x: x[0])
                slide_files = slide_files[:max_slides]

                for slide_num, slide_file in slide_files:
                    try:
                        content = zip_ref.read(slide_file).decode(
                            "utf-8", errors="ignore"
                        )
                        slide_text = self._extract_text_from_xml(content)
                        if slide_text:
                            texts.append(slide_text)
                    except:
                        continue

            return " ".join(texts)
        except Exception as e:
            print(f"快速提取失败: {str(e)}")
            # 回退到完整提取
            return self._extract_from_pptx_full(file_path, max_slides)

    def _extract_from_pptx_full(self, file_path: str, max_slides: int) -> str:
        """完整提取.pptx文本（使用python-pptx）"""
        try:
            from pptx import Presentation
            from pptx.shapes.base import BaseShape  # 导入BaseShape用于类型提示

            prs = Presentation(file_path)
            texts = []

            # 使用列表切片获取前max_slides个幻灯片
            slides = list(prs.slides)[:max_slides]

            for i, slide in enumerate(slides):
                slide_texts = []
                for shape in slide.shapes:
                    # 使用try-except处理可能的属性访问错误
                    try:
                        if hasattr(shape, "text") and shape.text:
                            text = shape.text.strip()
                            if text and len(text) > 1:
                                slide_texts.append(text)
                    except AttributeError:
                        continue

                if slide_texts:
                    texts.append(" ".join(slide_texts))

            return " ".join(texts)
        except ImportError:
            print("警告: 未安装python-pptx，无法解析PPT文件")
            return ""
        except Exception as e:
            print(f"PPT解析出错: {str(e)}")
            return ""

    def _extract_from_ppt_old(self, file_path: str, max_slides: int) -> str:
        """提取旧版.ppt文本（支持有限）"""
        warnings.warn("旧版.ppt文件支持有限，建议转换为.pptx格式")

        try:
            # 尝试使用python-pptx（也支持部分.ppt）
            return self._extract_from_pptx_full(file_path, max_slides)
        except:
            # 如果失败，返回空字符串
            return ""

    def _extract_text_from_xml(self, xml_content: str) -> str:
        """从XML内容提取文本"""
        # 简单正则提取
        patterns = [
            r"<a:t[^>]*>([^<]+)</a:t>",  # 文本元素
            r"<p:txBody>.*?<a:t>([^<]+)</a:t>",  # 文本框
        ]

        texts = []
        for pattern in patterns:
            matches = re.findall(pattern, xml_content, re.DOTALL)
            for match in matches:
                if isinstance(match, str):
                    text = match.strip()
                    if text:
                        texts.append(text)

        return " ".join(texts)

    def get_slide_count(self, file_path: str) -> int:
        """获取幻灯片数量"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            return len(prs.slides)
        except:
            return 0

    def extract_images(
        self, file_path: str, output_dir: str, max_images: int = 10
    ) -> List[str]:
        """提取PPT中的图像（训练阶段用）"""
        if not file_path.lower().endswith(".pptx"):
            return []

        image_paths = []
        try:
            from pptx import Presentation
            from pptx.shapes.base import BaseShape  # 导入BaseShape用于类型提示
            import os

            prs = Presentation(file_path)
            image_count = 0

            slides = list(prs.slides)
            for i, slide in enumerate(slides):
                if image_count >= max_images:
                    break

                for shape in slide.shapes:
                    if image_count >= max_images:
                        break

                    # 检查形状是否有image属性（使用try-except避免类型错误）
                    try:
                        # python-pptx中，只有图片形状才有image属性
                        if hasattr(shape, "image"):
                            image = shape.image
                            image_bytes = image.blob

                            # 确定图片格式
                            ext = image.ext
                            if not ext:
                                ext = ".png"

                            # 保存图片
                            image_filename = f"slide_{i}_img_{image_count}{ext}"
                            image_path = os.path.join(output_dir, image_filename)

                            with open(image_path, "wb") as f:
                                f.write(image_bytes)

                            image_paths.append(image_path)
                            image_count += 1
                    except AttributeError:
                        # 形状没有image属性，跳过
                        continue
                    except Exception as e:
                        print(f"处理图像时出错: {str(e)}")
                        continue

            return image_paths
        except Exception as e:
            print(f"提取图像失败: {str(e)}")
            return []
