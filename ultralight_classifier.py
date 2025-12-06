"""
UltraLight PPT Classifier - 轻量级示例模块

此模块实现了一个尽量简单、易部署的 `UltraLightPPTClassifier` 类。
在实际工程中可根据 `docs/optimization_plan.md` 中的格式替换或扩展。
"""

import os
import re
import gc
import zipfile
import joblib
import numpy as np

try:
    from pptx import Presentation
except Exception:
    Presentation = None


class UltraLightPPTClassifier:
    """极简PPT分类器（示例实现）

    仅作示例：加载 joblib 保存的压缩结构并提供 `predict_fast`。
    """

    def __init__(self, model_path):
        if not os.path.exists(model_path):
            raise FileNotFoundError(f"模型文件不存在: {model_path}")

        data = joblib.load(model_path)

        self.vectorizer = data.get("vectorizer")
        self.model = data.get("model")
        self.feature_mask = data.get("feature_mask")
        self.keyword_matcher = data.get("keyword_matcher", {})
        self.precomputed = data.get("precomputed", {})
        # 小优化：压缩词表并量化模型参数（在安全范围内）
        if hasattr(self.vectorizer, "vocabulary_"):
            try:
                self._optimize_memory()
            except Exception:
                # 忽略任何优化时的错误，保证加载成功
                pass

    def extract_text_fast(self, file_path, max_slides=30):
        """快速从 pptx 提取纯文本（优先解析 XML）"""
        try:
            texts = []
            with zipfile.ZipFile(file_path, "r") as zf:
                slide_files = [
                    f for f in zf.namelist() if f.startswith("ppt/slides/slide")
                ]
                for slide_file in slide_files[:max_slides]:
                    try:
                        content = zf.read(slide_file).decode("utf-8", errors="ignore")
                        matches = re.findall(r"<a:t[^>]*>([^<]+)</a:t>", content)
                        texts.extend(matches)
                    except Exception:
                        continue

            return " ".join(texts)[:5000]
        except Exception:
            return self._extract_text_fallback(file_path)

    def _extract_text_fallback(self, file_path):
        if Presentation is None:
            return ""

        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides[:20]:
                for shape in slide.shapes[:50]:
                    if hasattr(shape, "text"):
                        t = shape.text.strip()
                        if t:
                            texts.append(t)

            del prs
            gc.collect()
            return " ".join(texts)[:5000]
        except Exception:
            return ""

    def _optimize_memory(self):
        """内存优化：压缩词表并量化模型参数"""
        vocab = getattr(self.vectorizer, "vocabulary_", None)
        if vocab is None or self.feature_mask is None:
            return

        try:
            mask_set = set(int(i) for i in np.asarray(self.feature_mask).ravel())
        except Exception:
            mask_set = set()

        try:
            new_vocab = {k: v for k, v in vocab.items() if int(v) in mask_set}
            # 如果 new_vocab 为空则保留原 vocab
            if new_vocab:
                self.vectorizer.vocabulary_ = new_vocab
        except Exception:
            pass

        # 将模型参数量化为 float16（若存在）
        if hasattr(self.model, "coef_"):
            try:
                self.model.coef_ = self.model.coef_.astype(np.float16)
            except Exception:
                pass

        # 清理停用词属性以减少内存（如果存在）
        if hasattr(self.vectorizer, "stop_words_"):
            try:
                delattr(self.vectorizer, "stop_words_")
            except Exception:
                try:
                    del self.vectorizer.stop_words_
                except Exception:
                    pass

    def _keyword_fast_match(self, text):
        # 返回 (label, score)
        best = ("未知", 0.0)
        for label, keywords in (self.keyword_matcher or {}).items():
            if not keywords:
                continue
            count = 0
            for kw in keywords:
                if kw and kw in text:
                    count += 1
            if count > 0:
                score = count / max(1, len(keywords))
                if score > best[1]:
                    best = (label, min(score, 0.95))
        return best

    def _extract_light_features(self, text):
        if self.vectorizer is None or not hasattr(self.vectorizer, "transform"):
            return None

        text_clean = re.sub(r"\s+", " ", text.lower())
        try:
            features = self.vectorizer.transform([text_clean])
            if self.feature_mask is not None:
                features = features[:, self.feature_mask]

            # 对非常低维的情况返回稠密
            if features.shape[1] < 500:
                return features.toarray()
            return features
        except Exception:
            return None

    def predict_fast(self, ppt_file):
        text = self.extract_text_fast(ppt_file)
        if not text or len(text) < 8:
            return "未知", 0.0

        fast_label, fast_conf = self._keyword_fast_match(text)
        if fast_conf > 0.8:
            # 如果关键词匹配器返回的是索引字符串或数字，解析为可读标签
            try:
                idx = int(fast_label)
                return self._resolve_label(idx), fast_conf
            except Exception:
                return fast_label, fast_conf

        features = self._extract_light_features(text)
        if features is None:
            return fast_label, fast_conf

        try:
            if hasattr(self.model, "predict_proba"):
                proba = self.model.predict_proba(features)[0]
                idx = int(np.argmax(proba))
                confidence = float(proba[idx])
                label = self._resolve_label(idx)
            else:
                pred = self.model.predict(features)[0]
                label = str(pred)
                confidence = 0.7

            confidence = max(confidence, fast_conf * 0.3)
            return label, min(confidence, 0.99)
        except Exception:
            return fast_label, fast_conf

    def _resolve_label(self, idx):
        """解析 label_mapping，返回 index 对应的可读标签。"""
        lm = self.precomputed.get("label_mapping")
        # lm 可能为 list 或 dict
        if isinstance(lm, (list, tuple)):
            if 0 <= idx < len(lm):
                return lm[idx] or str(idx)
            return str(idx)

        if isinstance(lm, dict):
            # 如果是 dict 且键为索引（如 '0': '语文' 或 0: '语文'）
            try:
                key = idx
                if key in lm:
                    return lm[key]
            except Exception:
                pass

            # 如果 dict 为 label->index，尝试反转
            for k, v in lm.items():
                try:
                    if int(v) == idx:
                        return k
                except Exception:
                    continue

        return str(idx)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("model", help="path to deployment_model.joblib")
    parser.add_argument("ppt", nargs="?", help="pptx file to test", default=None)
    args = parser.parse_args()

    clf = UltraLightPPTClassifier(args.model)
    if args.ppt and os.path.exists(args.ppt):
        print(clf.predict_fast(args.ppt))
    else:
        print("Loaded model:", args.model)
