#!/usr/bin/env python3
"""
生成训练用PPT脚本 - 创建各学科高中知识点PPT
"""

import os
import random
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Tuple

# 确保数据目录存在
DATA_DIR = Path("data/raw")
for subject in ["语文", "数学", "英语", "物理", "化学", "生物"]:
    (DATA_DIR / subject).mkdir(parents=True, exist_ok=True)


class PPTGenerator:
    """PPT生成器"""

    def __init__(self):
        self.subjects = ["语文", "数学", "英语", "物理", "化学", "生物"]
        self.subject_colors = {
            "语文": RGBColor(66, 133, 244),  # 蓝色
            "数学": RGBColor(219, 68, 55),  # 红色
            "英语": RGBColor(244, 180, 0),  # 黄色
            "物理": RGBColor(15, 157, 88),  # 绿色
            "化学": RGBColor(156, 39, 176),  # 紫色
            "生物": RGBColor(0, 172, 193),  # 青色
        }

    def generate_subject_ppt(
        self, subject: str, topic: str, content: List[str], output_path: str
    ):
        """
        生成学科PPT

        Args:
            subject: 学科名称
            topic: 主题
            content: 内容列表
            output_path: 输出路径
        """
        # 创建演示文稿
        prs = Presentation()

        # 添加标题幻灯片
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = f"{subject} - {topic}"

        # 设置副标题
        subtitle = slide.placeholders[1]
        subtitle.text = f"高中知识点 | {datetime.now().strftime('%Y年%m月%d日')}"

        # 添加目录页
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "目录"

        # 添加目录内容
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()

        for i, item in enumerate(content, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {item}"
            p.level = 0
            p.font.size = Pt(24)

        # 添加内容页
        for i, item in enumerate(content, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"第{i}节: {item}"

            # 添加内容
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()

            # 根据学科添加具体内容
            detailed_content = self._generate_detailed_content(subject, item)
            for line in detailed_content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
                p.font.size = Pt(18)

            # 添加例子（随机添加）
            if random.random() > 0.3:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = f"例子: {item}"

                content_placeholder = slide.placeholders[1]
                tf = content_placeholder.text_frame
                tf.clear()

                examples = self._generate_examples(subject, item)
                for example in examples:
                    p = tf.add_paragraph()
                    p.text = example
                    p.level = 0
                    p.font.size = Pt(18)

        # 添加总结页
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "总结"

        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()

        summary = self._generate_summary(subject, topic)
        for line in summary:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(20)

        # 保存PPT
        prs.save(output_path)
        print(f"已生成: {output_path}")

    def _generate_detailed_content(self, subject: str, topic: str) -> List[str]:
        """生成详细内容"""
        content_map = {
            "语文": self._generate_chinese_content,
            "数学": self._generate_math_content,
            "英语": self._generate_english_content,
            "物理": self._generate_physics_content,
            "化学": self._generate_chemistry_content,
            "生物": self._generate_biology_content,
        }

        generator = content_map.get(subject, lambda x: [f"{subject}内容: {x}"])
        return generator(topic)

    def _generate_examples(self, subject: str, topic: str) -> List[str]:
        """生成例子"""
        examples_map = {
            "语文": self._generate_chinese_examples,
            "数学": self._generate_math_examples,
            "英语": self._generate_english_examples,
            "物理": self._generate_physics_examples,
            "化学": self._generate_chemistry_examples,
            "生物": self._generate_biology_examples,
        }

        generator = examples_map.get(subject, lambda x: [f"例子: 这是{x}的一个示例"])
        return generator(topic)

    def _generate_summary(self, subject: str, topic: str) -> List[str]:
        """生成总结"""
        return [
            f"本节学习了{topic}的相关知识",
            f"重点内容包括:",
            "• 概念理解",
            "• 应用方法",
            "• 典型例题",
            "课后建议: 完成相关练习题",
        ]

    def _generate_chinese_content(self, topic: str) -> List[str]:
        """生成语文内容"""
        contents = [
            f"{topic}是高中语文的重要知识点",
            "主要涉及以下几个方面:",
            "• 文学常识",
            "• 文言文阅读",
            "• 现代文阅读",
            "• 作文技巧",
        ]

        # 根据不同主题添加特定内容
        if "文言文" in topic:
            contents.extend(
                [
                    "文言文特点:",
                    "- 古汉语语法",
                    "- 常用虚词",
                    "- 古今异义",
                    "- 特殊句式",
                ]
            )
        elif "诗歌" in topic:
            contents.extend(
                [
                    "诗歌鉴赏要点:",
                    "- 意象分析",
                    "- 情感表达",
                    "- 艺术手法",
                    "- 时代背景",
                ]
            )

        return contents

    def _generate_math_content(self, topic: str) -> List[str]:
        """生成数学内容"""
        contents = [
            f"{topic}是高中数学的核心概念",
            "主要内容包括:",
            "• 定义与性质",
            "• 公式推导",
            "• 解题方法",
            "• 应用实例",
        ]

        if "函数" in topic:
            contents.extend(["函数三要素:", "- 定义域", "- 值域", "- 对应关系"])
        elif "几何" in topic:
            contents.extend(
                ["几何证明方法:", "- 综合法", "- 分析法", "- 反证法", "- 坐标法"]
            )

        return contents

    def _generate_english_content(self, topic: str) -> List[str]:
        """生成英语内容"""
        contents = [
            f"{topic} is an important topic in high school English",
            "Main contents include:",
            "• Grammar rules",
            "• Vocabulary building",
            "• Reading comprehension",
            "• Writing skills",
        ]

        if "grammar" in topic.lower() or "语法" in topic:
            contents.extend(
                [
                    "Key grammar points:",
                    "- Tense and aspect",
                    "- Sentence structure",
                    "- Parts of speech",
                    "- Clauses and phrases",
                ]
            )
        elif "reading" in topic.lower() or "阅读" in topic:
            contents.extend(
                [
                    "Reading strategies:",
                    "- Skimming and scanning",
                    "- Context clues",
                    "- Inference skills",
                    "- Critical thinking",
                ]
            )

        return contents

    def _generate_physics_content(self, topic: str) -> List[str]:
        """生成物理内容"""
        contents = [
            f"{topic}是高中物理的重要知识点",
            "主要内容包括:",
            "• 物理概念",
            "• 定律公式",
            "• 实验方法",
            "• 应用实例",
        ]

        if "力学" in topic:
            contents.extend(
                [
                    "力学基本定律:",
                    "- 牛顿三定律",
                    "- 动量守恒",
                    "- 能量守恒",
                    "- 万有引力",
                ]
            )
        elif "电" in topic:
            contents.extend(
                ["电学基本概念:", "- 电场强度", "- 电势差", "- 欧姆定律", "- 电磁感应"]
            )

        return contents

    def _generate_chemistry_content(self, topic: str) -> List[str]:
        """生成化学内容"""
        contents = [
            f"{topic}是高中化学的重要知识点",
            "主要内容包括:",
            "• 化学反应原理",
            "• 物质性质",
            "• 实验操作",
            "• 计算应用",
        ]

        if "反应" in topic:
            contents.extend(
                [
                    "化学反应类型:",
                    "- 氧化还原反应",
                    "- 酸碱反应",
                    "- 置换反应",
                    "- 化合分解反应",
                ]
            )
        elif "有机" in topic:
            contents.extend(
                [
                    "有机化学特点:",
                    "- 同分异构",
                    "- 官能团性质",
                    "- 反应机理",
                    "- 合成路线",
                ]
            )

        return contents

    def _generate_biology_content(self, topic: str) -> List[str]:
        """生成生物内容"""
        contents = [
            f"{topic}是高中生物的重要知识点",
            "主要内容包括:",
            "• 生命现象",
            "• 生理过程",
            "• 遗传规律",
            "• 生态关系",
        ]

        if "细胞" in topic:
            contents.extend(
                ["细胞结构:", "- 细胞膜", "- 细胞核", "- 细胞器", "- 细胞周期"]
            )
        elif "遗传" in topic:
            contents.extend(
                ["遗传规律:", "- 孟德尔定律", "- DNA结构", "- 基因表达", "- 遗传变异"]
            )

        return contents

    def _generate_chinese_examples(self, topic: str) -> List[str]:
        """生成语文例子"""
        examples = [
            f"{topic}的例子:",
            "例1: 文言文翻译",
            '"学而时习之，不亦说乎" - 《论语》',
            "翻译: 学习并且时常温习，不也是很愉快吗？",
            "",
            "例2: 诗歌鉴赏",
            "《静夜思》 李白",
            "床前明月光，疑是地上霜。",
            "举头望明月，低头思故乡。",
        ]

        if "作文" in topic:
            examples.extend(
                [
                    "作文写作技巧:",
                    "1. 立意要新颖",
                    "2. 结构要清晰",
                    "3. 语言要生动",
                    "4. 论证要充分",
                ]
            )

        return examples

    def _generate_math_examples(self, topic: str) -> List[str]:
        """生成数学例子"""
        examples = [
            f"{topic}的例子:",
            "例1: 函数求解",
            "已知 f(x) = 2x + 3, 求 f(5)",
            "解: f(5) = 2×5 + 3 = 13",
            "",
            "例2: 几何证明",
            "证明: 三角形内角和为180°",
        ]

        if "方程" in topic:
            examples.extend(
                [
                    "解方程例子:",
                    "2x² + 3x - 5 = 0",
                    "解: 使用求根公式",
                    "x = [-3 ± √(3² - 4×2×(-5))] / (2×2)",
                    "x = [-3 ± √49] / 4",
                    "x₁ = 1, x₂ = -2.5",
                ]
            )

        return examples

    def _generate_english_examples(self, topic: str) -> List[str]:
        """生成英语例子"""
        examples = [
            f"Examples for {topic}:",
            "",
            "Example 1: Sentence structure",
            "Simple sentence: I study English.",
            "Compound sentence: I study English, and he studies Math.",
            "",
            "Example 2: Vocabulary",
            "Important words:",
            "- analyze (v.) 分析",
            "- conclusion (n.) 结论",
            "- evidence (n.) 证据",
        ]

        if "writing" in topic.lower() or "写作" in topic:
            examples.extend(
                [
                    "Writing example:",
                    "Topic: My favorite book",
                    "Introduction: I would like to introduce my favorite book.",
                    "Body: It tells an inspiring story about friendship.",
                    "Conclusion: This book teaches me valuable life lessons.",
                ]
            )

        return examples

    def _generate_physics_examples(self, topic: str) -> List[str]:
        """生成物理例子"""
        examples = [
            f"{topic}的例子:",
            "例1: 牛顿第二定律",
            "F = ma",
            "其中: F为力, m为质量, a为加速度",
            "",
            "例2: 自由落体",
            "h = ½gt²",
            "g ≈ 9.8 m/s² (重力加速度)",
        ]

        if "电" in topic:
            examples.extend(
                [
                    "电路计算例子:",
                    "已知: U = 12V, R = 4Ω",
                    "求: 电流 I",
                    "解: I = U/R = 12/4 = 3A",
                ]
            )

        return examples

    def _generate_chemistry_examples(self, topic: str) -> List[str]:
        """生成化学例子"""
        examples = [
            f"{topic}的例子:",
            "例1: 化学反应",
            "2H₂ + O₂ → 2H₂O",
            "氢气与氧气反应生成水",
            "",
            "例2: 摩尔计算",
            "计算1 mol H₂O的质量:",
            "M(H₂O) = 2×1 + 16 = 18 g/mol",
        ]

        if "酸碱" in topic:
            examples.extend(
                [
                    "酸碱中和反应:",
                    "HCl + NaOH → NaCl + H₂O",
                    "盐酸与氢氧化钠反应生成氯化钠和水",
                ]
            )

        return examples

    def _generate_biology_examples(self, topic: str) -> List[str]:
        """生成生物例子"""
        examples = [
            f"{topic}的例子:",
            "例1: 细胞分裂",
            "有丝分裂过程:",
            "1. 前期 2. 中期",
            "3. 后期 4. 末期",
            "",
            "例2: 孟德尔遗传",
            "豌豆实验:",
            "高茎(显性) × 矮茎(隐性)",
            "F₁代全部为高茎",
        ]

        if "光合" in topic:
            examples.extend(
                ["光合作用方程式:", "6CO₂ + 6H₂O → C₆H₁₂O₆ + 6O₂", "光能转化为化学能"]
            )

        return examples


class TrainingPPTGenerator:
    """训练用PPT生成器"""

    def __init__(self):
        self.generator = PPTGenerator()
        self.subjects_topics = {
            "语文": [
                "文言文阅读技巧",
                "古代诗歌鉴赏",
                "现代文阅读分析",
                "作文写作指导",
                "语言文字应用",
                "文学常识积累",
                "修辞手法解析",
                "名著导读",
                "古诗文默写",
                "病句修改",
            ],
            "数学": [
                "函数与方程",
                "三角函数",
                "立体几何",
                "解析几何",
                "概率统计",
                "数列与极限",
                "导数与应用",
                "向量运算",
                "不等式证明",
                "复数概念",
            ],
            "英语": [
                "英语语法精讲",
                "词汇记忆方法",
                "阅读理解技巧",
                "写作训练指导",
                "听力提高策略",
                "口语表达练习",
                "完形填空解题",
                "短文改错方法",
                "翻译技巧",
                "英语文化背景",
            ],
            "物理": [
                "力学基础",
                "运动学规律",
                "牛顿定律应用",
                "功与能量",
                "电路分析",
                "电磁感应",
                "光学原理",
                "热学基础",
                "原子物理",
                "物理实验方法",
            ],
            "化学": [
                "化学反应原理",
                "化学平衡",
                "有机化学基础",
                "无机化学",
                "化学实验操作",
                "溶液计算",
                "电化学",
                "化学键与分子",
                "元素周期律",
                "化学与生活",
            ],
            "生物": [
                "细胞结构与功能",
                "遗传与进化",
                "生态系统",
                "人体生理",
                "生物技术",
                "植物生理",
                "微生物学",
                "生物进化论",
                "实验生物学",
                "生物与环境",
            ],
        }

    def generate_all(self, samples_per_subject: int = 10):
        """为每个学科生成多个PPT"""
        print("开始生成训练用PPT...")
        print("=" * 50)

        for subject, topics in self.subjects_topics.items():
            print(f"\n生成 {subject} PPT...")
            subject_dir = DATA_DIR / subject

            # 确保目录存在
            subject_dir.mkdir(exist_ok=True)

            # 生成指定数量的PPT
            selected_topics = (
                topics[:samples_per_subject]
                if len(topics) >= samples_per_subject
                else topics
            )

            for i, topic in enumerate(selected_topics, 1):
                # 生成内容大纲
                content = self._generate_content_outline(subject, topic)

                # 生成文件名
                filename = f"{subject}_高中{i:02d}_{topic}.pptx"
                output_path = subject_dir / filename

                # 生成PPT
                try:
                    self.generator.generate_subject_ppt(
                        subject=subject,
                        topic=topic,
                        content=content,
                        output_path=str(output_path),
                    )
                except Exception as e:
                    print(f"  生成失败 {filename}: {str(e)}")

            print(f"  {subject}: 已生成 {len(selected_topics)} 个PPT")

    def _generate_content_outline(self, subject: str, topic: str) -> List[str]:
        """生成内容大纲"""
        outlines = {
            "语文": [
                f"{topic}概述",
                "基础知识讲解",
                "重点难点分析",
                "典型例题解析",
                "学习方法指导",
                "课后练习",
            ],
            "数学": [
                f"{topic}定义",
                "性质与定理",
                "公式推导",
                "解题方法",
                "应用实例",
                "练习题",
            ],
            "英语": [
                f"{topic} Introduction",
                "Key Concepts",
                "Usage Rules",
                "Practice Exercises",
                "Common Mistakes",
                "Review Questions",
            ],
            "物理": [
                f"{topic}概念",
                "物理定律",
                "公式应用",
                "实验方法",
                "问题解析",
                "知识拓展",
            ],
            "化学": [
                f"{topic}原理",
                "化学反应",
                "物质性质",
                "实验操作",
                "计算应用",
                "知识总结",
            ],
            "生物": [
                f"{topic}概述",
                "生命现象",
                "生理过程",
                "实验观察",
                "应用实例",
                "知识延伸",
            ],
        }

        return outlines.get(
            subject, [f"{topic}简介", "主要内容", "重点难点", "应用实例", "总结练习"]
        )


def main():
    """主函数"""
    import argparse

    parser = argparse.ArgumentParser(description="生成训练用各学科PPT")
    parser.add_argument(
        "--samples", type=int, default=10, help="每个学科生成的PPT数量 (默认: 10)"
    )
    parser.add_argument(
        "--output", type=str, default="data/raw", help="输出目录 (默认: data/raw)"
    )

    args = parser.parse_args()

    # 更新数据目录
    global DATA_DIR
    DATA_DIR = Path(args.output)

    # 生成PPT
    generator = TrainingPPTGenerator()
    generator.generate_all(samples_per_subject=args.samples)

    print("\n" + "=" * 50)
    print("PPT生成完成!")
    print("=" * 50)

    # 统计生成的PPT数量
    total_count = 0
    for subject_dir in DATA_DIR.iterdir():
        if subject_dir.is_dir():
            ppt_files = list(subject_dir.glob("*.pptx"))
            count = len(ppt_files)
            print(f"{subject_dir.name}: {count} 个PPT")
            total_count += count

    print(f"\n总计: {total_count} 个PPT")
    print(f"保存在: {DATA_DIR}/学科名/ 目录下")


if __name__ == "__main__":
    main()
