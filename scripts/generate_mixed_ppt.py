#!/usr/bin/env python3
"""
生成混合内容PPT - 用于测试分类器的鲁棒性
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import random
from pathlib import Path


def generate_mixed_ppt():
    """生成混合学科内容的PPT"""

    # 学科关键词
    subject_keywords = {
        "语文": ["文言文", "诗歌", "散文", "修辞", "作文", "鲁迅", "杜甫", "红楼梦"],
        "数学": ["函数", "方程", "几何", "概率", "导数", "积分", "矩阵", "证明"],
        "英语": ["grammar", "vocabulary", "reading", "writing", "tense", "essay"],
        "物理": ["力", "加速度", "电场", "磁场", "电路", "能量", "动量", "量子"],
        "化学": ["元素", "分子", "反应", "化学式", "溶液", "酸碱", "有机"],
        "生物": ["细胞", "基因", "DNA", "蛋白质", "生态系统", "进化", "遗传"],
    }

    # 创建混合内容
    mixed_contents = []
    all_subjects = list(subject_keywords.keys())

    # 随机选择3-4个学科混合
    selected_subjects = random.sample(all_subjects, random.randint(3, 4))

    for subject in selected_subjects:
        keywords = subject_keywords[subject]
        # 随机选择2-3个关键词
        selected_keywords = random.sample(keywords, random.randint(2, 3))
        mixed_contents.append(f"{subject}相关内容: {', '.join(selected_keywords)}")

    # 生成PPT
    prs = Presentation()

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "综合复习资料"
    subtitle = slide.placeholders[1]
    subtitle.text = "多学科知识点整理"

    # 混合内容页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "多学科知识点"

    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()

    for content in mixed_contents:
        p = tf.add_paragraph()
        p.text = "• " + content
        p.level = 0
        p.font.size = Pt(20)

    # 保存
    output_dir = Path("data/test/mixed")
    output_dir.mkdir(parents=True, exist_ok=True)

    output_path = output_dir / f"混合内容_{random.randint(1000, 9999)}.pptx"
    prs.save(output_path)

    print(f"生成混合内容PPT: {output_path}")
    print("包含学科:", ", ".join(selected_subjects))

    return selected_subjects, str(output_path)


def generate_mixed_dataset(num_ppts: int = 10):
    """生成多个混合内容PPT"""
    print(f"生成 {num_ppts} 个混合内容PPT...")
    print("=" * 50)

    results = []
    for i in range(num_ppts):
        print(f"\n生成第 {i+1}/{num_ppts} 个...")
        subjects, path = generate_mixed_ppt()
        results.append(
            {
                "file": path,
                "subjects": subjects,
                "primary_subject": subjects[0] if subjects else "未知",
            }
        )

    # 保存结果信息
    import json

    info_path = Path("data/test/mixed/info.json")
    with open(info_path, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)

    print(f"\n完成! 结果信息保存在: {info_path}")
    return results


if __name__ == "__main__":
    generate_mixed_dataset(5)
